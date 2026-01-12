# presentation/generate_ppt.py
# Builds a 15-slide PowerPoint for "Project Management Evaluation System"
# Output: ../exports/PMES_Presentation.pptx

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent
EXPORTS = ROOT / "exports"
EXPORTS.mkdir(parents=True, exist_ok=True)
OUT = EXPORTS / "PMES_Presentation.pptx"

prs = Presentation()

# Helpers
TITLE = 0
TITLE_AND_CONTENT = 1
TITLE_ONLY = 5

def add_title(text, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE])
    slide.shapes.title.text = text
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullets(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, (list, tuple)):
            # First element is the heading, rest are sub-bullets
            heading = item[0]
            p.text = heading
            p.font.size = Pt(18)
            p.level = 0
            for sub in item[1:]:
                sp = tf.add_paragraph()
                sp.text = str(sub)
                sp.level = 1
                sp.font.size = Pt(16)
        else:
            p.text = str(item)
            p.level = 0
            p.font.size = Pt(18)

# 1. Title
add_title(
    "Project Management Evaluation System",
    "CIE Review‑1 — Department / Team / Date"
)

# 2–3. Introduction / Problem
add_bullets("Introduction", [
    "Multi‑evaluator assessment of final‑year projects",
    "Manual collation is error‑prone; rubric needs standardization",
    "Target rubric: 50 marks (Lit 20, Problem 10, Presentation 10, Q&A 10)",
    "Need per‑student sheets and consolidated exports for audits",
])
add_bullets("Problem Statement & Scope", [
    "Web app to ingest .xlsx, validate, compute components/averages",
    "Persist Student/Evaluation data; export CSV/Excel",
    "Scope: Review‑1 flow; per‑student + consolidated outputs",
    "Out of scope now: auth/roles, multi‑review timeline, analytics",
])

# 4–8. Literature Review (5 slides)
add_bullets("Literature Review — Rubrics", [
    "Measurable criteria and weightages improve fairness",
    "Inter‑rater reliability considerations",
    "Standardization increases transparency",
    "(Insert 2–3 references)",
])
add_bullets("Literature Review — Aggregation", [
    "Mean vs median vs trimmed mean across evaluators",
    "Handling disagreement and outliers",
    "Mapping totals to components (proportional splits)",
    "(Insert 2–3 references)",
])
add_bullets("Literature Review — Ingestion & Integrity", [
    ".xlsx preferred over CSV for structure and safety",
    "Header normalization and schema validation",
    "User feedback for invalid cells",
    "(Insert 2 references)",
])
add_bullets("Literature Review — Data Modeling", [
    "Simple ER: Student 1—* Evaluation",
    "Uniqueness: one evaluation per student per review",
    "SQLite for simplicity; MySQL for multi‑user",
    "(Insert 2 references)",
])
add_bullets("Literature Review — Reporting", [
    "Standardized outputs for audits/NAAC/NBA",
    "Automated templates improve traceability",
    "(Insert 1–2 references)",
])

# 9. Summary of Literature
add_bullets("Summary of Literature Review", [
    "Standard rubrics + reliable multi‑evaluator aggregation",
    "Robust .xlsx ingestion and clear validation",
    "Canonical storage with constraints and reproducible exports",
])

# 10. Objectives
add_bullets("Objectives", [
    ["Functional:", "Upload standardized .xlsx", "Replace database with upload", "Compute components/averages", "Export CSV/Excel"],
    ["Quality:", "Data correctness; idempotent imports", "Simple UI for viva", "Audit‑ready artifacts"],
])

# 11–12. Design
add_bullets("Design — Architecture", [
    ["Client:", "Upload, Students list, Detail, Downloads"],
    ["Server (Flask):", "/upload, /students, detail, downloads"],
    ["Data:", "SQLAlchemy models: Student, Evaluation; SQLite/MySQL"],
    ["Files:", "templates/, exports/, excel_template.py"],
])
add_bullets("Design — Data & Flow", [
    ["ER:", "Student(id, seat_no, name, group_no, project_title)", "Evaluation(review_no, components, per‑evaluator, total, student_id)", "Unique(student_id, review_no)"],
    ["Ingestion:", "Read header → normalize → map keys", "3 evaluator totals → average + proportional components", "Or components → sum", "Or total → reverse‑engineer split (Hamilton rounding)"],
])

# 13–14. Implementation
add_bullets("Implementation — Key Modules", [
    "app.py: routes, .xlsx parsing, replace‑import, exports",
    "utils.py: normalize_header, reverse_engineer_components, TOTAL_MAX=50",
    "models.py: SQLAlchemy models + unique constraint",
    "excel_template.py: per‑student Review‑1 workbook",
    "Validations: required columns; numeric casting; flash errors",
])
add_bullets("Implementation — Demo & Outputs", [
    "Upload: select .xlsx → replaces existing data → success",
    "Students: table with per‑component and totals; download links",
    "Detail: evaluator sums, component breakdown, Excel download",
    "Exports: exports/evaluations_review1.csv (deterministic)",
])

# 15. Conclusion and Future Scope
add_bullets("Conclusion & Future Scope", [
    ["Conclusion:", "Standardized Review‑1 with multi‑evaluator support", ".xlsx ingestion & audit‑ready exports", "Reliable and fast during viva"],
    ["Future:", "Review‑2/Final timeline", "Auth/roles", "PDF/email automation", "Analytics & dashboards", "Cloud deployment & backups"],
])

prs.save(str(OUT))
print(f"Created: {OUT}")
