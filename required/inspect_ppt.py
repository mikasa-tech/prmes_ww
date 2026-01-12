from pptx import Presentation
import sys
from pathlib import Path

def dump_pptx(file_path):
    path = Path(file_path)
    if not path.exists():
        print(f"File not found: {file_path}")
        return

    try:
        prs = Presentation(path)
        print(f"--- Successfully opened: {path.name} ---")
        print(f"Total Slides: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        print(f"  - {text}")
                
                if shape.has_table:
                    print("  [Table Found]")
                    for row in shape.table.rows:
                        row_text = [cell.text_frame.text.strip() for cell in row.cells]
                        print(f"    | {' | '.join(row_text)} |")
    except Exception as e:
        print(f"Error reading PPTX: {e}")

if __name__ == "__main__":
    target = r"C:\Users\asus\OneDrive\Documents\prmes5 (1)3.pptx"
    dump_pptx(target)
